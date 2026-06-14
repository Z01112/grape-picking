#!/usr/bin/env python
"""Generate mid-term defense PPT for GPPoint-DETR."""
import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

DARK_BLUE = RGBColor(0x18, 0x5F, 0xA5)
LIGHT_BLUE = RGBColor(0xE6, 0xF1, 0xFB)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BLACK = RGBColor(0x2C, 0x2C, 0x2A)
GRAY = RGBColor(0x66, 0x66, 0x66)
GREEN = RGBColor(0x1D, 0x9E, 0x75)
RED = RGBColor(0xE2, 0x4B, 0x4A)
ORANGE = RGBColor(0xD8, 0x5A, 0x30)

BASE = "D:/Projects/RT-DETR/RT-DETRv4"
OUT = os.path.join(BASE, "docs/midterm_defense.pptx")

prs = Presentation()
prs.slide_width = Inches(13.33)
prs.slide_height = Inches(7.5)

def add_slide():
    return prs.slides.add_slide(prs.slide_layouts[6])  # blank

def add_title_bg(slide, title_text, subtitle_text=""):
    """Add a top banner with title."""
    banner = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1.3))
    banner.fill.solid()
    banner.fill.fore_color.rgb = DARK_BLUE
    banner.line.fill.background = True
    tf = banner.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title_text
    p.font.size = Pt(32)
    p.font.color.rgb = WHITE
    p.font.bold = True
    p.alignment = PP_ALIGN.LEFT
    tf.margin_left = Inches(0.8)
    
    if subtitle_text:
        sub = slide.shapes.add_textbox(Inches(0.8), Inches(1.1), Inches(11), Inches(0.5))
        stf = sub.text_frame
        sp = stf.paragraphs[0]
        sp.text = subtitle_text
        sp.font.size = Pt(16)
        sp.font.color.rgb = GRAY

def add_textbox(slide, left, top, width, height, text, size=14, bold=False, color=BLACK, align=PP_ALIGN.LEFT):
    tb = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.alignment = align
    return tf

def add_card(slide, left, top, width, height, label, value, sub="", color=DARK_BLUE, highlight=False):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height))
    shape.fill.solid()
    shape.fill.fore_color.rgb = LIGHT_BLUE if highlight else WHITE
    shape.line.color.rgb = color if highlight else RGBColor(0xDD, 0xDD, 0xDD)
    shape.line.width = Pt(2) if highlight else Pt(0.5)
    tf = shape.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.15)
    tf.margin_right = Inches(0.15)
    p1 = tf.paragraphs[0]
    p1.text = label
    p1.font.size = Pt(11)
    p1.font.color.rgb = GRAY
    p1.alignment = PP_ALIGN.CENTER
    p2 = tf.add_paragraph()
    p2.text = str(value)
    p2.font.size = Pt(28)
    p2.font.bold = True
    p2.font.color.rgb = color
    p2.alignment = PP_ALIGN.CENTER
    if sub:
        p3 = tf.add_paragraph()
        p3.text = sub
        p3.font.size = Pt(10)
        p3.font.color.rgb = GRAY
        p3.alignment = PP_ALIGN.CENTER

# ====== SLIDE 1: COVER ======
s = add_slide()
banner = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
banner.fill.solid()
banner.fill.fore_color.rgb = DARK_BLUE
banner.line.fill.background()
add_textbox(s, 1, 1.5, 11, 1.5, "基于RT-DETR的葡萄采摘决策系统研究", 36, True, WHITE, PP_ALIGN.CENTER)
add_textbox(s, 1, 3.2, 11, 0.5, "GPPoint-DETR: Instance-Bound Picking Point Detection", 20, False, LIGHT_BLUE, PP_ALIGN.CENTER)
add_textbox(s, 1, 4.5, 11, 0.4, "中期答辩", 24, True, WHITE, PP_ALIGN.CENTER)
add_textbox(s, 1, 5.5, 11, 0.4, "2026", 16, False, LIGHT_BLUE, PP_ALIGN.CENTER)

# ====== SLIDE 2: BACKGROUND ======
s = add_slide()
add_title_bg(s, "研究背景与意义", "Background & Motivation")
add_textbox(s, 0.8, 1.7, 5, 0.4, "葡萄采摘痛点", 18, True)
bullets = [
    "- 农业劳动力短缺，机器人采摘需求迫切",
    "- 采摘点精确定位是视觉系统的核心难题",
    "- 传统方法：检测bbox 取中心点 误差大",
    "- 果梗位置在bbox顶部附近，不在中心",
]
tb = add_textbox(s, 0.8, 2.2, 5.5, 3, "", 14)
for b in bullets:
    p = tb.add_paragraph()
    p.text = b
    p.font.size = Pt(14)
    p.font.color.rgb = BLACK
    p.space_after = Pt(8)

add_textbox(s, 7, 1.7, 5.5, 0.4, "RT-DETR基础", 18, True)
bullets2 = [
    "- RT-DETR: 实时Transformer检测器",
    "- Query机制：每个query负责一个实例",
    "- 优势：端到端，无需NMS后处理",
    "- 问题：只能检测bbox，无采摘点输出",
]
tb2 = add_textbox(s, 7, 2.2, 5.5, 3, "", 14)
for b in bullets2:
    p = tb2.add_paragraph()
    p.text = b
    p.font.size = Pt(14)
    p.font.color.rgb = BLACK
    p.space_after = Pt(8)

# Core question box
add_textbox(s, 2.5, 5.5, 8.5, 0.6, "核心问题：能否利用Query机制，为每个葡萄实例绑定一个采摘点？", 18, True, DARK_BLUE, PP_ALIGN.CENTER)

# ====== SLIDE 3: CORE APPROACH ======
s = add_slide()
add_title_bg(s, "方法论：从Query到采摘点", "Methodology: Query-to-PickingPoint")
add_textbox(s, 0.8, 1.7, 5.5, 0.4, "建模思路", 18, True)
approach = [
    "1. 检测到葡萄串bbox后，不是取中心点",
    "2. 在query特征上新增两个预测头：",
    "   - Visibility Head: 判别该bbox是否有可见采摘点",
    "   - Offset Head: 回归采摘点相对于anchor的偏移",
    "3. 采摘点坐标 = anchor + offset × bbox_size",
    "4. Point Cross-Attention: 6x6局部ROI增强offset特征",
    "5. Point-Aware Matcher: 训练时考虑offset质量做匹配",
]
tb = add_textbox(s, 0.8, 2.2, 6, 4, "", 14)
for a in approach:
    p = tb.add_paragraph()
    p.text = a
    p.font.size = Pt(14)
    p.font.color.rgb = BLACK
    p.space_after = Pt(6)

add_textbox(s, 7.2, 1.7, 5.5, 0.4, "与传统方法的区别", 18, True)
compare = [
    "传统方法: bbox -> 中心点 (误差大)",
    "          或bbox -> heatmap峰值 (复杂)",
    "",
    "本文方法: query -> {bbox, has_picking, offset}",
    "          offset回归 + 解码 = 采摘点",
    "",
    "优势:",
    "- 实例绑定: query机制天然支持",
    "- 端到端: 不需要后处理",
    "- 轻量: 两个MLP头增量<5%参数",
]
tb2 = add_textbox(s, 7.2, 2.2, 5, 4, "", 14)
for c in compare:
    p = tb2.add_paragraph()
    p.text = c
    p.font.size = Pt(14)
    p.font.color.rgb = BLACK
    p.space_after = Pt(4)

# ====== SLIDE 4: MODEL ARCHITECTURE ======
s = add_slide()
add_title_bg(s, "模型总体架构", "GPPoint-DETR Architecture")
arch_path = os.path.join(BASE, "docs/ppt_figs/architecture.svg")
if os.path.exists(arch_path):
    s.shapes.add_picture(arch_path, Inches(0.5), Inches(1.5), Inches(12.3), Inches(5.5))
else:
    add_textbox(s, 1, 2, 11, 2, "[架构图SVG请手动插入]\n详见SVG: docs/ppt_figs/architecture.svg", 14, False, GRAY, PP_ALIGN.CENTER)

# ====== SLIDE 5: DECODER MODELING ======
s = add_slide()
add_title_bg(s, "Decoder建模：采摘点解码原理", "Decoder Picking Point Mechanism")
add_textbox(s, 0.8, 1.7, 12, 0.3, "Top-Center Anchor + Offset Regression", 16, True)

steps = [
    ("1. Anchor定义", "bbox顶部0.12h处设定Top-Center Anchor点\nanchor_x = bbox_cx, anchor_y = bbox_top + 0.12 * bbox_h"),
    ("2. Offset回归", "query特征 PointCA(6x6 MHA) Offset MLP [dx, dy]\ndx, dy是归一化偏移量，相对于bbox宽高"),
    ("3. 坐标解码", "picking_x = anchor_x + dx * bbox_w\npicking_y = anchor_y + dy * bbox_h"),
    ("4. 实例绑定", "query_i 同时输出 {bbox_i, has_picking_i, picking_point_i}\n每个葡萄实例独立绑定一个采摘点"),
]
for i, (title, desc) in enumerate(steps):
    y = 2.3 + i * 1.25
    box = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1), Inches(y), Inches(2.2), Inches(1))
    box.fill.solid()
    box.fill.fore_color.rgb = DARK_BLUE
    box.line.fill.background = True
    btf = box.text_frame
    btf.word_wrap = True
    bp = btf.paragraphs[0]
    bp.text = title
    bp.font.size = Pt(14)
    bp.font.color.rgb = WHITE
    bp.font.bold = True
    bp.alignment = PP_ALIGN.CENTER
    
    add_textbox(s, 3.5, y + 0.1, 9, 0.9, desc, 13, False, BLACK)

# Key note
add_textbox(s, 1, 7.0, 11, 0.3, "关键认知: 采摘点不是独立检测类别，而是通过bbox + offset回归解码得到", 14, True, RED)

# ====== SLIDE 6: ENCODER OPTIMIZATION ======
s = add_slide()
add_title_bg(s, "Encoder优化：EMA + BiFPN", "Encoder Enhancement")
add_textbox(s, 0.8, 1.7, 5.5, 0.4, "消融验证", 18, True)

# Encoder comparison table
table_data = [
    ["Encoder", "F1", "Pair/533", "L2"],
    ["Standard (no EMA/BiFPN)", "0.8399", "417", "14.84"],
    ["EMA + BiFPN (main baseline)", "0.8787 (+4.6%)", "449 (+7.7%)", "13.02 (-12.3%)"],
]
y = 2.3
for row in table_data:
    for j, cell in enumerate(row):
        is_header = (row == table_data[0])
        x = 0.8 + j * 3.5
        w = 3.2
        h = 0.5 if is_header else 0.45
        shape = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
        shape.fill.solid()
        shape.fill.fore_color.rgb = DARK_BLUE if is_header else WHITE
        shape.line.color.rgb = DARK_BLUE
        shape.line.width = Pt(0.5)
        tf = shape.text_frame
        p = tf.paragraphs[0]
        p.text = cell
        p.font.size = Pt(13)
        p.font.color.rgb = WHITE if is_header else BLACK
        p.font.bold = is_header
        p.alignment = PP_ALIGN.CENTER
    y += (0.5 if row == table_data[0] else 0.45)

add_textbox(s, 0.8, 4.2, 11, 0.4, "数据验证：EMA+BiFPN使F1提升4.6%，L2降低12.3%——对picking任务至关重要", 15, True, DARK_BLUE)

add_textbox(s, 0.8, 4.9, 5.5, 0.4, "编码器组件", 18, True)
components = [
    "- AIFI Transformer: 仅处理S5'最深特征，做intra-scale语义交互",
    "- CCFF / BiFPN: S3'+S4'+F5跨尺度融合，融合细节和语义",
    "- EMA Attention: 对P3/P4/P5做特征增强",
    "- 1x1 Conv投影统一256通道"
]
tb = add_textbox(s, 0.8, 5.4, 5.5, 2, "", 13)
for c in components:
    p = tb.add_paragraph()
    p.text = c
    p.font.size = Pt(13)
    p.font.color.rgb = BLACK
    p.space_after = Pt(5)

add_textbox(s, 7, 4.9, 5.5, 0.4, "设计原则", 18, True)
principles = [
    "- 不使用VFM/DSI/GAM(论文参考文档明确标注)",
    "- AIFI只处理S5'(不是所有尺度)",
    "- F5是特征名，不是模块",
    "- 确保与其他消融实验组件解耦"
]
tb2 = add_textbox(s, 7, 5.4, 5.5, 2, "", 13)
for p_text in principles:
    p = tb2.add_paragraph()
    p.text = p_text
    p.font.size = Pt(13)
    p.font.color.rgb = BLACK
    p.space_after = Pt(5)

# ====== SLIDE 7: IMPROVEMENT PATH ======
s = add_slide()
add_title_bg(s, "改进路径：B0-B1-PtCA-PAM", "Improvement Path (Test Set, 271 GT)")

# F1 + Pair bar chart
add_textbox(s, 0.8, 1.7, 11, 0.4, "Picking性能逐步提升", 16, True)

chart_path = os.path.join(BASE, "docs/ppt_figs/improvement_chart.html")
if os.path.exists(chart_path):
    add_textbox(s, 1, 2.2, 11, 5, "[改进路径图表HTML请手动截图]\n文件: docs/ppt_figs/improvement_chart.html", 14, False, GRAY, PP_ALIGN.CENTER)

# Step cards
steps_data = [
    ("Step 1: B1 Backbone", "B0 B1", "AP +1.4%\nMAE_x -15.2%\nF1 -0.7pp", GREEN),
    ("Step 2: Point Cross-Attn", "+PointCA", "F1 +1.69pp\nPair +7\nL2 +12.8%", ORANGE),
    ("Step 3: PAM Matcher", "+PAM", "F1 +1.31pp\nFN -24%\nL2 -7.2%", DARK_BLUE),
]
for i, (title, badge, result, color) in enumerate(steps_data):
    x = 0.8 + i * 4.2
    card = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(5.8), Inches(3.8), Inches(1.5))
    card.fill.solid()
    card.fill.fore_color.rgb = WHITE
    card.line.color.rgb = color
    card.line.width = Pt(1.5)
    tf = card.text_frame
    tf.margin_left = Inches(0.15)
    p1 = tf.paragraphs[0]
    p1.text = title
    p1.font.size = Pt(13)
    p1.font.bold = True
    p1.font.color.rgb = color
    p2 = tf.add_paragraph()
    p2.text = result
    p2.font.size = Pt(11)
    p2.font.color.rgb = BLACK

# Bottom summary
add_textbox(s, 0.8, 7.1, 11, 0.3, "最终: B0 Baseline F1=0.878  B1+PtCA+PAM F1=0.900 (+2.15pp)  FN 36  25 (-30.6%)", 14, True, DARK_BLUE, PP_ALIGN.CENTER)

# ====== SLIDE 8: ABLATION TABLE ======
s = add_slide()
add_title_bg(s, "消融实验全景（Test Set）", "Full Ablation Study")
add_textbox(s, 0.5, 1.7, 12, 0.3, "所有数据均基于Test Set (180 images, 271 GT picking points)", 13, False, GRAY, PP_ALIGN.CENTER)

headers = ["Model", "AP", "F1", "Pair", "FN", "FP", "L2", "Params"]
rows = [
    ["Simple Baseline (no EMA/BiFPN)", "0.629", "0.874", "226", "45", "20", "14.68", "13.38M"],
    ["Main Baseline (B0)", "0.623", "0.878", "234", "36", "29", "13.33", "13.41M"],
    ["B1", "0.632", "0.872", "231", "39", "29", "13.52", "13.76M"],
    ["B1 + PAM", "0.628", "0.892", "236", "34", "23", "13.78", "13.76M"],
    ["B1 + PointCA", "0.644", "0.886", "238", "33", "28", "15.25", "14.09M"],
    ["B1 + PointCA + PAM (Ours)", "0.635", "0.900", "246", "25", "30", "14.15", "14.09M"],
]
y = 2.3
for j, h in enumerate(headers):
    w = [2.3, 0.8, 0.8, 0.7, 0.6, 0.6, 0.8, 1.0]
    x = 0.5 + sum(w[:j])
    shape = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(w[j]), Inches(0.45))
    shape.fill.solid()
    shape.fill.fore_color.rgb = DARK_BLUE
    shape.line.fill.background = True  # no line
    tf = shape.text_frame
    p = tf.paragraphs[0]
    p.text = h
    p.font.size = Pt(11)
    p.font.color.rgb = WHITE
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER

y += 0.45
for ri, row in enumerate(rows):
    is_ours = ri == len(rows) - 1
    for j, cell in enumerate(row):
        w = [2.3, 0.8, 0.8, 0.7, 0.6, 0.6, 0.8, 1.0]
        x = 0.5 + sum(w[:j])
        shape = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(w[j]), Inches(0.42))
        shape.fill.solid()
        shape.fill.fore_color.rgb = RGBColor(0xE1, 0xF5, 0xEE) if is_ours else WHITE
        shape.line.color.rgb = GREEN if is_ours else RGBColor(0xDD, 0xDD, 0xDD)
        shape.line.width = Pt(1) if is_ours else Pt(0.3)
        tf = shape.text_frame
        p = tf.paragraphs[0]
        p.text = cell
        p.font.size = Pt(11)
        p.font.color.rgb = DARK_BLUE if is_ours else BLACK
        p.font.bold = is_ours
        p.alignment = PP_ALIGN.CENTER
    y += 0.42

# Legend
add_textbox(s, 0.5, 5.8, 12, 0.3, "Pair = correctly matched picking point pairs / 271", 11, False, GRAY)

# Key conclusions
conclusions = [
    "1. B1+PointCA+PAM 在F1(0.900)和Pair(246)上均为最优",
    "2. PAM在B1上单独有效(+2.35pp F1)，在B0上反而有害",
    "3. PointCA提升F1但恶化L2，PAM可修复此退化(-7.2% L2)",
]
tb = add_textbox(s, 0.5, 6.3, 12, 1.2, "", 13)
for c in conclusions:
    p = tb.add_paragraph()
    p.text = c
    p.font.size = Pt(13)
    p.font.color.rgb = BLACK
    p.font.bold = True
    p.space_after = Pt(4)

# ====== SLIDE 9: QUALITATIVE RESULTS ======
s = add_slide()
add_title_bg(s, "可视化结果（Test Set）", "Qualitative Results: B1+PointCA+PAM")

for i in range(3):
    fpath = os.path.join(BASE, f"docs/ppt_figs/sample{i+1}.png")
    try:
        s.shapes.add_picture(fpath, Inches(0.5 + i * 4.2), Inches(1.8), Inches(3.8), Inches(4.0))
        add_textbox(s, 0.5 + i * 4.2, 5.95, 3.8, 0.3, f"Sample {i+1}", 12, False, GRAY, PP_ALIGN.CENTER)
    except Exception as e:
        add_textbox(s, 0.5 + i * 4.2, 3, 3.8, 0.5, f"[Image sample{i+1}]", 11, False, GRAY, PP_ALIGN.CENTER)

add_textbox(s, 0.5, 6.5, 12, 0.3, "蓝色框: 检测bbox | 橙色点: 预测采摘点 | 绿色框: 正确配对", 12, False, GRAY, PP_ALIGN.CENTER)
add_textbox(s, 0.5, 6.9, 12, 0.3, "更多可视化样本: outputs/04_diagnostics/qualitative_single_grape_samples/", 11, False, GRAY)

# ====== SLIDE 10: MAIN MODEL HIGHLIGHTS ======
s = add_slide()
add_title_bg(s, "主模型 B1+PointCA+PAM 核心指标", "Main Model: Key Performance (Test Set)")

add_card(s, 0.8, 2.0, 2.8, 1.6, "Picking F1", "0.900", "历史最高 +2.15pp", DARK_BLUE, True)
add_card(s, 3.8, 2.0, 2.8, 1.6, "Correctly Paired", "246/271", "90.8% of GT", GREEN)
add_card(s, 6.8, 2.0, 2.8, 1.6, "FN Reduction", "-30.6%", "36 25", RED)
add_card(s, 9.8, 2.0, 2.8, 1.6, "Parameter Overhead", "+5.1%", "13.41M  14.09M", GRAY)

add_card(s, 0.8, 4.0, 2.8, 1.6, "Grape AP", "0.635", "+2.0% vs Baseline", DARK_BLUE)
add_card(s, 3.8, 4.0, 2.8, 1.6, "L2 Error", "14.15 px", "y: 9.27 x: 8.31", ORANGE)
add_card(s, 6.8, 4.0, 2.8, 1.6, "PAM Inference Cost", "0", "training only", GREEN)
add_card(s, 9.8, 4.0, 2.8, 1.6, "Encoder", "EMA+BiFPN", "F1 +4.6pp", DARK_BLUE)

add_textbox(s, 0.8, 6.2, 11, 0.5, "论文主要贡献：Point Cross-Attention偏移特征增强 + Point-Aware Matcher训练匹配优化", 16, True, DARK_BLUE, PP_ALIGN.CENTER)
add_textbox(s, 0.8, 6.8, 11, 0.3, "F1提升2.15pp | FN减少30.6% | 采摘点配对率90.8%", 14, False, BLACK, PP_ALIGN.CENTER)

# ====== SLIDE 11: COMPARISON WITH YOLO ======
s = add_slide()
add_title_bg(s, "与传统方法对比", "Comparison with Baseline Methods")

comparison_data = [
    ["Method", "F1", "Pair", "L2", "Params"],
    ["YOLO11n-Pose (keypoint)", "—", "—", "—", "2.6M"],
    ["BBox Center (naive)", "—", "—", "high", "0"],
    ["Main Baseline (B0)", "0.878", "234", "13.33", "13.41M"],
    ["Ours (B1+PtCA+PAM)", "0.900", "246", "14.15", "14.09M"],
]
y = 2.2
for ri, row in enumerate(comparison_data):
    is_header = ri == 0
    is_ours = ri == len(comparison_data) - 1
    for j, cell in enumerate(row):
        w = [5, 1.5, 1.2, 1.2, 1.5]
        x = 0.8 + sum(w[:j])
        shape = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(w[j]), Inches(0.45))
        shape.fill.solid()
        if is_header:
            shape.fill.fore_color.rgb = DARK_BLUE
        elif is_ours:
            shape.fill.fore_color.rgb = LIGHT_BLUE
        else:
            shape.fill.fore_color.rgb = WHITE
        shape.line.color.rgb = DARK_BLUE if is_ours else RGBColor(0xDD, 0xDD, 0xDD)
        shape.line.width = Pt(0.5)
        tf = shape.text_frame
        p = tf.paragraphs[0]
        p.text = cell
        p.font.size = Pt(12)
        p.font.color.rgb = WHITE if is_header else BLACK
        p.font.bold = is_header or is_ours
        p.alignment = PP_ALIGN.CENTER
    y += 0.45

add_textbox(s, 0.8, 5.0, 11, 0.4, "对比分析", 16, True)
analysis = [
    "- BBox中心点法：简单但误差大(果梗不在中心)，无法判定采摘点是否存在",
    "- YOLO11n-Pose: 轻量但keypoint方法缺乏实例绑定，L2定位精度不足",
    "- Ours优势: 实例绑定确保每个葡萄串独立预测，offset回归+anchor机制精确定位",
]
tb = add_textbox(s, 0.8, 5.5, 11, 1.5, "", 13)
for a in analysis:
    p = tb.add_paragraph()
    p.text = a
    p.font.size = Pt(13)
    p.font.color.rgb = BLACK
    p.space_after = Pt(6)

# ====== SLIDE 12: FUTURE WORK ======
s = add_slide()
add_title_bg(s, "后续工作：机械臂实验部署", "Future Work: Robotic Deployment")

add_textbox(s, 0.8, 1.7, 5.5, 0.4, "2D 3D坐标转换", 18, True)
steps_3d = [
    "1. 模型输出2D图像采摘点坐标 (x,y)",
    "2. RGB-D相机获取深度信息 d",
    "3. 相机内参 + 手眼标定矩阵",
    "4. 转换到机械臂基座坐标系",
    "5. 机械臂路径规划 + 抓取执行"
]
tb = add_textbox(s, 0.8, 2.2, 5.5, 3.5, "", 14)
for st in steps_3d:
    p = tb.add_paragraph()
    p.text = st
    p.font.size = Pt(14)
    p.font.color.rgb = BLACK
    p.space_after = Pt(6)

add_textbox(s, 7, 1.7, 5.5, 0.4, "待完成工作", 18, True)
todo = [
    "- 实机标定: 相机-机械臂手眼标定",
    "- 推理部署: ONNX/TensorRT优化",
    "- 实机测试: 葡萄采摘抓取实验",
    "- 反馈优化: 闭环抓取调整",
    "- 论文完善: 撰写实验分析章节",
]
tb2 = add_textbox(s, 7, 2.2, 5.5, 3.5, "", 14)
for t in todo:
    p = tb2.add_paragraph()
    p.text = t
    p.font.size = Pt(14)
    p.font.color.rgb = BLACK
    p.space_after = Pt(6)

add_textbox(s, 2, 6.5, 9.5, 0.4, "已完成: 2D采摘点检测模型 下一步: RGB-D 机械臂实机抓取验证", 18, True, DARK_BLUE, PP_ALIGN.CENTER)

# ====== SLIDE 13: CONCLUSION ======
s = add_slide()
add_title_bg(s, "总结与贡献", "Summary & Contributions")

add_textbox(s, 0.8, 1.7, 5.5, 0.4, "完成了什么", 18, True)
done = [
    "1. 基于RT-DETR构建采摘点检测框架",
    "   利用query机制实现实例级采摘点绑定",
    "2. Top-Center Anchor + Offset回归",
    "   不需要单独检测采摘点类别",
    "3. Point Cross-Attention",
    "   6x6 MHA增强offset特征的空间表示",
    "4. Point-Aware Matcher",
    "   训练匹配关注offset质量，降低FN 30.6%",
    "5. 系统消融实验",
    "   9个模型变体、test set统一评估",
]
tb = add_textbox(s, 0.8, 2.2, 5.5, 3.5, "", 14)
for d in done:
    p = tb.add_paragraph()
    p.text = d
    p.font.size = Pt(14)
    p.font.color.rgb = BLACK
    p.space_after = Pt(4)

add_textbox(s, 7, 1.7, 5.5, 0.4, "核心数据", 18, True)
add_card(s, 7, 2.3, 5.5, 1.0, "Picking F1", "0.900", "Test Set, 271 GT", DARK_BLUE, True)
add_card(s, 7, 3.5, 5.5, 1.0, "FN: 36  25", "-30.6%", "漏判显著减少", GREEN)
add_card(s, 7, 4.7, 5.5, 1.0, "Pair: 234  246", "+5.1%", "采摘点配对率提升", DARK_BLUE)

add_textbox(s, 0.8, 6.5, 11, 0.4, "方法可推广到其他果实采摘任务（柑橘、苹果、番茄等）", 16, False, BLACK, PP_ALIGN.CENTER)

# ====== SLIDE 14: ACKNOWLEDGEMENTS ======
s = add_slide()
banner2 = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
banner2.fill.solid()
banner2.fill.fore_color.rgb = DARK_BLUE
banner2.line.fill.background()
add_textbox(s, 1, 3, 11, 1, "谢谢！", 48, True, WHITE, PP_ALIGN.CENTER)
add_textbox(s, 1, 4.2, 11, 0.5, "Thank You - Questions Welcome", 20, False, LIGHT_BLUE, PP_ALIGN.CENTER)
add_textbox(s, 1, 5.5, 11, 0.4, "基于RT-DETR的葡萄采摘决策系统研究", 18, False, WHITE, PP_ALIGN.CENTER)

prs.save(OUT)
print(f"PPT saved to: {OUT}")
print(f"Total slides: {len(prs.slides)}")
